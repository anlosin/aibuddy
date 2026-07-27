"""PowerPoint 演示文稿插件 — 创建/读取 .pptx 文件（对应 WorkBuddy 的 pptx 文件处理能力）"""
import os

PLUGIN_INFO = {
    "name": "pptx",
    "description": "创建 PowerPoint 演示文稿(.pptx)或读取已有文稿内容。支持封面、多张幻灯片，每页可含标题、正文与要点列表。",
    "version": "1.0",
}

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "create_pptx",
            "description": "创建一个 PowerPoint 演示文稿(.pptx)。提供标题（用作封面）和若干幻灯片，每页可含标题、正文与要点列表。文件保存在程序同级目录下。",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "文件名，如 presentation.pptx"
                    },
                    "title": {
                        "type": "string",
                        "description": "演示文稿标题（显示在封面页）"
                    },
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {
                                    "type": "string",
                                    "description": "幻灯片标题"
                                },
                                "content": {
                                    "type": "string",
                                    "description": "幻灯片正文段落内容"
                                },
                                "bullets": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "要点列表（可选），每项为一行要点"
                                }
                            },
                            "required": ["heading"]
                        },
                        "description": "幻灯片列表，每页包含 heading(标题)、content(正文)、bullets(可选要点列表)"
                    }
                },
                "required": ["filename", "title"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_pptx",
            "description": "读取一个 PowerPoint 演示文稿的内容，返回每页的标题与文本。",
            "parameters": {
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "PPTX 文件路径"
                    }
                },
                "required": ["filepath"]
            }
        }
    },
]


def _safe_path(filepath):
    if os.path.isabs(filepath):
        return filepath
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", filepath)


def execute(name, arguments):
    if name == "create_pptx":
        return _create(arguments)
    if name == "read_pptx":
        return _read(arguments)
    return f"未知工具: {name}"


def _create(args):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "错误: 请先安装 python-pptx (pip install python-pptx)"

    filename = args.get("filename", "presentation.pptx")
    title = args.get("title", "未命名演示文稿")
    slides = args.get("slides", [])

    prs = Presentation()

    # 封面页（标题 + 副标题布局）
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "由 AI 助手生成"

    if not slides:
        prs.slides.add_slide(prs.slide_layouts[1])
    else:
        for s in slides:
            heading = s.get("heading", "")
            content = s.get("content", "")
            bullets = s.get("bullets", []) or []
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            first = True
            if content:
                p = tf.paragraphs[0]
                p.text = content
                first = False
            for b in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = b
                first = False

    path = _safe_path(filename)
    prs.save(path)
    return (
        f"PowerPoint 演示文稿已创建: {filename}\n"
        f"路径: {os.path.abspath(path)}\n"
        f"幻灯片数: {len(prs.slides)}"
    )


def _read(args):
    try:
        from pptx import Presentation
    except ImportError:
        return "错误: 请先安装 python-pptx (pip install python-pptx)"

    filepath = args.get("filepath", "")
    if not filepath:
        return "错误: 未提供文件路径"

    path = _safe_path(filepath)
    if not os.path.exists(path):
        return f"错误: 文件不存在 - {path}"

    try:
        prs = Presentation(path)
    except Exception as e:
        return f"无法打开文件: {e}"

    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- 第 {i} 页 ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)

    content = "\n".join(lines)
    if len(content) > 6000:
        content = content[:6000] + "\n...[内容已截断]"

    return f"[{filepath}]\n页数: {len(prs.slides)}\n\n{content}"
